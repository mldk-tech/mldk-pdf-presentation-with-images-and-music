import os
import time # Added in case we want a delay for debugging, not currently in use
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

# --- Settings ---
IMAGES_DIR = 'images'
MUSIC_DIR = 'music'
# Save as a Show file that will open directly in presentation mode. Change to 'presentation.pptx' for debugging opening issues.
OUTPUT_FILENAME = 'presentation.pptx'
ALLOWED_IMAGE_EXTENSIONS = ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif')
ALLOWED_MUSIC_EXTENSIONS = ('.mp3', '.wav', '.wma') # Can add more formats
SLIDE_TRANSITION_SECONDS = 1 # Automatic transition time in seconds

# --- Helper Functions ---

def set_automatic_transition(slide, seconds):
    """Sets an automatic transition for the slide after a number of seconds."""
    try:
        transition = slide.slide_show_transition
        transition.advance_on_click = False # Disable transition on click
        transition.advance_on_time = True   # Enable transition by time
        transition.advance_time = seconds * 1000 # Time in milliseconds
        # Can add a transition effect type here if desired, for example:
        # from pptx.enum.transitions import PP_TRANSITION_EFFECT
        # transition.transition_effect = PP_TRANSITION_EFFECT.FADE
    except Exception as e:
        print(f"  Warning: Could not set transition for slide: {e}")


def add_image_slide(prs, img_path):
    """
    Adds a new slide with a centered image at maximum size while preserving aspect ratio,
    and sets an automatic transition.
    """
    print(f"-> Processing image: {img_path}")
    try:
        slide_layout = prs.slide_layouts[6]  # Use blank layout
        slide = prs.slides.add_slide(slide_layout)
        print(f"  Slide added for {os.path.basename(img_path)}")

        # Get slide dimensions in EMU
        slide_width_emu = prs.slide_width
        slide_height_emu = prs.slide_height

        # Add the image
        print(f"  Attempting to add picture object: {os.path.basename(img_path)}")
        pic = slide.shapes.add_picture(img_path, Emu(0), Emu(0), width=Emu(914400)) # Temporary initial size
        print(f"  Picture object added. Calculating optimal size...")

        # Get original image dimensions in pixels
        img_native_width_px = pic.image.size[0]
        img_native_height_px = pic.image.size[1]

        # Handle case of zero-size image (rare, but possible with a corrupted file)
        if img_native_width_px == 0 or img_native_height_px == 0:
             print(f"  Error: Image {os.path.basename(img_path)} has zero width or height. Skipping resize/positioning.")
             # Can decide to delete the slide or leave it blank
             # slide.shapes._spTree.remove(pic._element) # Remove the image itself if desired
             set_automatic_transition(slide, SLIDE_TRANSITION_SECONDS) # Still set transition
             return # Skip to continue

        # Calculate aspect ratio of the original image
        img_aspect_ratio = img_native_width_px / img_native_height_px

        # Calculate the maximum possible size within the slide
        if (slide_width_emu / slide_height_emu) > img_aspect_ratio:
            new_height_emu = slide_height_emu
            new_width_emu = int(new_height_emu * img_aspect_ratio)
        else:
            new_width_emu = slide_width_emu
            new_height_emu = int(new_width_emu / img_aspect_ratio)

        # Update image size on the slide
        pic.width = new_width_emu
        pic.height = new_height_emu
        print(f"  Resized picture to {new_width_emu}x{new_height_emu} EMU")

        # Calculate position (left, top) for centering
        pic.left = int((slide_width_emu - new_width_emu) / 2)
        pic.top = int((slide_height_emu - new_height_emu) / 2)
        print(f"  Positioned picture at {pic.left}, {pic.top} EMU")

        # Set automatic transition for this slide
        print(f"  Setting automatic transition ({SLIDE_TRANSITION_SECONDS}s)")
        set_automatic_transition(slide, SLIDE_TRANSITION_SECONDS)

        print(f"  Successfully processed image slide: {os.path.basename(img_path)}")

    except FileNotFoundError:
        print(f"  ERROR: Image file not found: {img_path}")
    except Exception as e:
        # Print a more detailed error
        import traceback
        print(f"  ERROR processing image {img_path}: {e}")
        # print(traceback.format_exc()) # Can uncomment to get full error details


def add_title_slide(prs, title_text):
    """
    Adds a title slide with given text and sets an automatic transition.
    """
    print(f"-> Adding title slide: {title_text}")
    try:
        title_slide_layout = prs.slide_layouts[5] # Title Only layout
        slide = prs.slides.add_slide(title_slide_layout)
        print(f"  Slide added for title '{title_text}'")

        title_shape = None
        # Try to find the title box defined in the layout
        if slide.shapes.title:
            title_shape = slide.shapes.title
            print("  Found title placeholder in layout.")
        # If not found, try to find any text placeholder (less ideal)
        elif slide.placeholders:
            for shape in slide.placeholders:
                # Look for the first placeholder that is a title or main body text
                # Placeholder types: TITLE=1, BODY=2, CENTER_TITLE=3, SUBTITLE=4 ...
                if shape.placeholder_format.idx in [1, 3]: # Title or Center Title
                    title_shape = shape
                    print(f"  Found placeholder (type {shape.placeholder_format.idx}) for title.")
                    break
                elif shape.placeholder_format.idx == 2: # Body text placeholder
                    title_shape = shape
                    print(f"  Found body placeholder (type {shape.placeholder_format.idx}) for title.")
                    # Don't break, prefer a title placeholder if found later

        if title_shape:
            title_shape.text = title_text
            # Basic formatting (optional)
            title_shape.text_frame.paragraphs[0].font.size = Pt(44)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE # Fit text size to shape
            print(f"  Set text for title '{title_text}'")
        else:
            # If no suitable placeholder is found, add a manual textbox
            print("  Warning: No suitable title placeholder found. Adding manual textbox.")
            left = top = Inches(0.5)
            width = prs.slide_width - Inches(1.0)
            height = Inches(1.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = title_text
            # Basic formatting for the manual box
            p = tf.paragraphs[0]
            p.font.size = Pt(40)
            p.font.bold = True
            tf.word_wrap = True
            # Can center the text in the box
            # from pptx.enum.text import PP_ALIGN
            # p.alignment = PP_ALIGN.CENTER
            # tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            print(f"  Added manual textbox for title '{title_text}'")

        # Set automatic transition for this slide
        print(f"  Setting automatic transition ({SLIDE_TRANSITION_SECONDS}s)")
        set_automatic_transition(slide, SLIDE_TRANSITION_SECONDS)

        print(f"  Successfully processed title slide: {title_text}")

    except Exception as e:
        import traceback
        print(f"  ERROR adding title slide for '{title_text}': {e}")
        # print(traceback.format_exc()) # Can uncomment to get full error details


def find_music_file(directory):
    """Searches for a supported audio file in the given directory."""
    if not os.path.isdir(directory):
        print(f"Music directory '{directory}' not found.")
        return None

    print(f"Searching for music file in '{directory}'...")
    for filename in os.listdir(directory):
        if filename.lower().endswith(ALLOWED_MUSIC_EXTENSIONS):
            found_path = os.path.join(directory, filename)
            print(f"  Found music file: {filename}")
            return found_path
    print("  No supported music file found.")
    return None

def add_background_music(prs, music_file_path):
    """
    Embeds a music file in the first slide.
    Note: Advanced settings require manual intervention or XML editing.
    """
    print(f"-> Attempting to add background music: {os.path.basename(music_file_path)}")
    if not prs.slides:
        print("  Cannot add music, no slides exist in the presentation.")
        return

    if not music_file_path or not os.path.exists(music_file_path):
        print(f"  ERROR: Music file path is invalid or file does not exist: '{music_file_path}'")
        return

    try:
        first_slide = prs.slides[0]
        print("  Adding music to the first slide.")

        # Determine MIME type (simplified)
        mime_type = 'audio/mpeg' # Default MP3
        if music_file_path.lower().endswith('.wav'):
            mime_type = 'audio/wav'
        elif music_file_path.lower().endswith('.wma'):
             mime_type = 'audio/x-ms-wma'
        print(f"  Using MIME type: {mime_type}")

        # Add the object - small and inconspicuous
        left = Inches(0.1)
        top = Inches(0.1)
        width = Inches(0.2) # Even smaller
        height = Inches(0.2)

        movie = first_slide.shapes.add_movie(
            music_file_path, left, top, width, height,
            mime_type=mime_type,
            poster_frame_image=None
        )
        print(f"  Successfully embedded music file '{os.path.basename(music_file_path)}'.")
        print("  Reminder: Manual setup in PowerPoint needed for 'Play Across Slides' and 'Loop'.")

    except Exception as e:
        import traceback
        print(f"  ERROR adding music file {music_file_path}: {e}")
        # print(traceback.format_exc())

# --- Main Script Section ---

if __name__ == "__main__":
    print("========================================")
    print("Starting presentation creation process...")
    print(f"Image directory: '{IMAGES_DIR}'")
    print(f"Music directory: '{MUSIC_DIR}'")
    print(f"Output file: '{OUTPUT_FILENAME}'")
    print(f"Slide transition time: {SLIDE_TRANSITION_SECONDS} second(s)")
    print("========================================")

    prs = Presentation()
    # Set default aspect ratio (16:9). Can change to 4:3 if needed.
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    print(f"Presentation object created with slide size: {prs.slide_width / Emu(12700):.2f} x {prs.slide_height / Emu(12700):.2f} inches")

    # Check if the images directory exists
    if not os.path.isdir(IMAGES_DIR):
        print(f"FATAL ERROR: Images directory '{IMAGES_DIR}' not found. Exiting.")
        exit(1) # Exit with error code

    items_in_images_dir = []
    try:
         # Sort items in the images directory to maintain alphabetical order
        items_in_images_dir = sorted(os.listdir(IMAGES_DIR))
        print(f"Found {len(items_in_images_dir)} items in '{IMAGES_DIR}'.")
    except Exception as e:
        print(f"FATAL ERROR: Could not list items in '{IMAGES_DIR}': {e}. Check permissions. Exiting.")
        exit(1)

    processed_slides = 0 # Counts how many slides were added

    # Iterate over items in the main images directory
    for item_name in items_in_images_dir:
        item_path = os.path.join(IMAGES_DIR, item_name)

        # If it's an image file in the main directory
        if os.path.isfile(item_path) and item_name.lower().endswith(ALLOWED_IMAGE_EXTENSIONS):
            add_image_slide(prs, item_path)
            processed_slides += 1

        # If it's a subdirectory
        elif os.path.isdir(item_path):
            print(f"\n--- Processing subdirectory: {item_name} ---")
            # Add a title slide with the subdirectory name
            add_title_slide(prs, item_name)
            processed_slides += 1

            sub_dir_items = []
            try:
                sub_dir_items = sorted(os.listdir(item_path))
                print(f"Found {len(sub_dir_items)} items in subdirectory '{item_name}'.")
            except Exception as e:
                print(f"  ERROR: Could not list items in subdirectory '{item_path}': {e}. Skipping.")
                continue # Skip to the next subdirectory

            # Iterate over images within the subdirectory
            for sub_item_name in sub_dir_items:
                sub_item_path = os.path.join(item_path, sub_item_name)
                if os.path.isfile(sub_item_path) and sub_item_name.lower().endswith(ALLOWED_IMAGE_EXTENSIONS):
                    add_image_slide(prs, sub_item_path)
                    processed_slides += 1
            print(f"--- Finished processing subdirectory: {item_name} ---\n")
        else:
             print(f"-> Skipping item (not a supported image or directory): {item_name}")


    # Add music if it exists
    print("--- Checking for background music ---")
    music_file = find_music_file(MUSIC_DIR)
    if music_file:
        if processed_slides > 0: # Only if there are slides to add music to
             add_background_music(prs, music_file)
        else:
             print("Skipping music addition as no slides were created.")
    # No message if the directory doesn't exist, because it was already printed in find_music_file

    # Save the presentation
    print("\n--- Saving presentation ---")
    if processed_slides > 0:
        try:
            prs.save(OUTPUT_FILENAME)
            print(f"Presentation saved successfully as '{OUTPUT_FILENAME}' ({processed_slides} slides created).")
            print("Reminder: Manual setup in PowerPoint might be needed for continuous presentation loop and cross-slide audio looping.")
            print(f"Tip: If '{OUTPUT_FILENAME}' fails to open, try changing the script to save as '.pptx' for debugging.")
        except Exception as e:
            import traceback
            print(f"FATAL ERROR saving presentation: {e}")
            # print(traceback.format_exc())
            print("The presentation file might be corrupted or incomplete.")
    else:
        print("No slides were added to the presentation. File not saved.")

    print("\n========================================")
    print("Script finished.")
    print("========================================")