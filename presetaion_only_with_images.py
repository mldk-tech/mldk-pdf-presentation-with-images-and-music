import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# --- Settings ---
IMAGE_DIR = 'images' # Name of the main directory containing images and subdirectories
OUTPUT_FILENAME = 'presentation_from_images.pptx' # Name of the presentation file to be created
SUPPORTED_EXTENSIONS = ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff') # Supported file extensions

# --- Helper Functions ---

def add_image_slide(prs, image_path):
    """
    Adds a new slide with a centered image at maximum size while maintaining aspect ratio.
    """
    try:
        # Use a blank slide layout (usually index 6)
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Get slide dimensions (in EMU - English Metric Units)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Add the picture but don't set dimensions yet, so we can get its original size
        # Temporarily place it at the corner
        pic = slide.shapes.add_picture(image_path, Inches(0), Inches(0))

        # Get the original aspect ratio of the image
        img_native_width = pic.image.size[0]
        img_native_height = pic.image.size[1]
        aspect_ratio = img_native_width / img_native_height

        # Calculate the maximum possible size while maintaining aspect ratio
        # Check if width limit or height limit is the critical one
        if (slide_width / slide_height) > aspect_ratio:
            # Height is limiting, adjust height to slide height
            height = slide_height
            width = aspect_ratio * height
        else:
            # Width is limiting, adjust width to slide width
            width = slide_width
            height = width / aspect_ratio

        # Calculate position to center the image on the slide
        left = (slide_width - width) / 2
        top = (slide_height - height) / 2

        # Update the image's position and size
        pic.left = int(left)
        pic.top = int(top)
        pic.width = int(width)
        pic.height = int(height)

        print(f"  - Added image: {os.path.basename(image_path)}")
        return True
    except Exception as e:
        print(f"  - Error adding image {os.path.basename(image_path)}: {e}")
        # Try to remove the shape if it was partially created and an error occurred
        try:
            sp = pic._element
            sp.getparent().remove(sp)
        except:
            pass # If we couldn't remove it, it's okay
        return False


def add_title_slide(prs, title_text):
    """
    Adds a title slide with the given text.
    """
    try:
        # Use a title slide layout (usually index 0 or 5)
        # Try 5 (Title Only) or 0 (Title Slide)
        try:
            title_slide_layout = prs.slide_layouts[5] # Try "Title Only"
        except IndexError:
            title_slide_layout = prs.slide_layouts[0] # If 5 doesn't exist, try "Title Slide"

        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        if title:
             title.text = title_text
             # Can adjust formatting if desired
             # title.text_frame.paragraphs[0].font.size = Pt(44)
             # title.text_frame.paragraphs[0].font.bold = True
        else:
            # If the layout doesn't have a 'title' shape, add a text box
             left = top = width = height = Inches(1.0)
             txBox = slide.shapes.add_textbox(left, top, width, height)
             tf = txBox.text_frame
             tf.text = title_text
             tf.paragraphs[0].font.size = Pt(40)
             tf.paragraphs[0].font.bold = True

        print(f"- Added title slide for: {title_text}")
    except Exception as e:
        print(f"- Error adding title slide for '{title_text}': {e}")

# --- Main Logic ---

def create_presentation():
    """
    Creates the presentation by scanning the images directory.
    """
    if not os.path.isdir(IMAGE_DIR):
        print(f"Error: Directory '{IMAGE_DIR}' not found.")
        print("Please create the 'images' directory and place your image files inside.")
        return

    # Create a new presentation object
    prs = Presentation()
    # Can set slide size if desired (e.g., 16:9), default is 4:3
    # prs.slide_width = Inches(16)
    # prs.slide_height = Inches(9)

    print(f"Scanning directory: '{IMAGE_DIR}'...")

    items_in_root = sorted(os.listdir(IMAGE_DIR)) # Sort items for consistent order

    # First, process images directly in the root directory
    print("Processing images in root directory...")
    root_images_processed = 0
    for item_name in items_in_root:
        item_path = os.path.join(IMAGE_DIR, item_name)
        if os.path.isfile(item_path) and item_name.lower().endswith(SUPPORTED_EXTENSIONS):
            if add_image_slide(prs, item_path):
                root_images_processed += 1

    if root_images_processed > 0:
         print(f"Processed {root_images_processed} images from the root directory.")
    else:
         print("No images found directly in the root directory.")


    # Process subdirectories
    print("\nProcessing subdirectories...")
    subdirs_processed = 0
    for item_name in items_in_root:
        item_path = os.path.join(IMAGE_DIR, item_name)
        if os.path.isdir(item_path):
            print(f"\nProcessing subdirectory: '{item_name}'")
            subdir_path = item_path
            subdir_name = item_name # The folder name will be used as the title

            # Add a title slide with the subdirectory name
            add_title_slide(prs, subdir_name)

            # Process images inside the subdirectory
            images_in_subdir = 0
            try:
                subdir_items = sorted(os.listdir(subdir_path)) # Sort here too
                for sub_item_name in subdir_items:
                    sub_item_path = os.path.join(subdir_path, sub_item_name)
                    if os.path.isfile(sub_item_path) and sub_item_name.lower().endswith(SUPPORTED_EXTENSIONS):
                         if add_image_slide(prs, sub_item_path):
                             images_in_subdir += 1
                if images_in_subdir > 0:
                     print(f"Processed {images_in_subdir} images from subdirectory '{item_name}'.")
                     subdirs_processed += 1
                else:
                    print(f"No images found in subdirectory '{item_name}'.")

            except Exception as e:
                print(f"Error processing subdirectory '{item_name}': {e}")


    if subdirs_processed > 0 :
        print(f"\nProcessed {subdirs_processed} subdirectories.")
    else:
        print("\nNo subdirectories with images found.")


    # Save the presentation
    try:
        prs.save(OUTPUT_FILENAME)
        print(f"\nPresentation saved successfully as '{OUTPUT_FILENAME}'")
    except Exception as e:
        print(f"\nError saving presentation: {e}")
        print("Please ensure you have write permissions in the current directory and the file is not open elsewhere.")

# --- Run the script ---
if __name__ == "__main__":
    create_presentation()